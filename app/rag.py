"""
RAG (Retrieval-Augmented Generation) core logic for NotesGPT.
Handles document loading, chunking, indexing, and querying.
Supports multiple document types: PDF, images (OCR), Word, Excel, PowerPoint, text, markdown.
"""

import os
import io
from typing import List, Dict, Any
from pathlib import Path

# LangChain imports
from langchain_community.document_loaders import PyPDFLoader, TextLoader, UnstructuredMarkdownLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_groq import ChatGroq
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_chroma import Chroma
from langchain_core.documents import Document

# Document processing imports
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    import pandas as pd
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from PIL import Image
    import pytesseract
    IMAGE_OCR_AVAILABLE = True
except ImportError:
    IMAGE_OCR_AVAILABLE = False



# ==================== Configuration Constants ====================
# Chunking parameters
CHUNK_SIZE = 800
CHUNK_OVERLAP = 150

# Query parameters
TOP_K = 5

# Model configuration
CHAT_MODEL_NAME = "llama-3.1-8b-instant"  # Free Groq model
EMBEDDING_MODEL_NAME = "sentence-transformers/all-MiniLM-L6-v2"  # Free HuggingFace embeddings

# Vector store configuration
CHROMA_DIR = "./chroma_db"
COLLECTION_NAME = "notesgpt"


# ==================== Model Initialization ====================
def get_embeddings():
    """Initialize and return HuggingFace embeddings model (free, local)."""
    return HuggingFaceEmbeddings(
        model_name=EMBEDDING_MODEL_NAME,
        model_kwargs={'device': 'cpu'},
        encode_kwargs={'normalize_embeddings': True}
    )


def get_llm():
    """Initialize and return ChatGroq model (free API)."""
    return ChatGroq(model=CHAT_MODEL_NAME, temperature=0)


def get_vector_store():
    """Get or create the Chroma vector store."""
    embeddings = get_embeddings()
    
    # Ensure directory exists
    Path(CHROMA_DIR).mkdir(parents=True, exist_ok=True)
    
    vector_store = Chroma(
        collection_name=COLLECTION_NAME,
        embedding_function=embeddings,
        persist_directory=CHROMA_DIR
    )
    
    return vector_store


# ==================== Document Loading ====================
def load_file_to_documents(path: str) -> List[Document]:
    """
    Load a file and convert it to LangChain Document objects.
    Supports: PDF, images (OCR), Word, Excel, PowerPoint, text, markdown.
    
    Args:
        path: Path to the file to load
        
    Returns:
        List of Document objects
    """
    file_extension = Path(path).suffix.lower()
    file_path = Path(path)
    
    try:
        # PDF files
        if file_extension == '.pdf':
            loader = PyPDFLoader(path)
            documents = loader.load()
            
            # Check if PDF contains any text
            total_text = "".join([doc.page_content.strip() for doc in documents])
            if not total_text or len(total_text) < 10:
                # Try OCR as fallback for scanned PDFs
                return [Document(
                    page_content=f"This appears to be a scanned PDF (image-based) with no extractable text. Please use OCR tools or re-upload as images for text extraction.",
                    metadata={"source": file_path.name, "error": True, "error_type": "scanned_pdf"}
                )]
        
        # Image files with OCR
        elif file_extension in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif']:
            documents = load_image_with_ocr(path)
        
        # Word documents
        elif file_extension in ['.docx', '.doc']:
            documents = load_word_document(path)
        
        # PowerPoint presentations
        elif file_extension in ['.pptx', '.ppt']:
            documents = load_powerpoint(path)
        
        # Excel spreadsheets
        elif file_extension in ['.xlsx', '.xls', '.csv']:
            documents = load_excel(path)
        
        # Markdown files
        elif file_extension in ['.md', '.markdown']:
            try:
                loader = UnstructuredMarkdownLoader(path)
                documents = loader.load()
            except:
                # Fallback to text loader
                loader = TextLoader(path, encoding='utf-8')
                documents = loader.load()
        
        # Plain text files
        else:
            loader = TextLoader(path, encoding='utf-8')
            documents = loader.load()
        
        return documents
    
    except Exception as e:
        # If all else fails, return an error document
        return [Document(
            page_content=f"Error loading file: {str(e)}",
            metadata={"source": file_path.name, "error": True}
        )]


def load_image_with_ocr(path: str) -> List[Document]:
    """
    Load an image and extract text using OCR.
    
    Args:
        path: Path to the image file
    
    Returns:
        List containing a single Document with extracted text
    """
    if not IMAGE_OCR_AVAILABLE:
        return [Document(
            page_content="Image processing not available. Please install Pillow and pytesseract.",
            metadata={"source": Path(path).name, "error": True}
        )]
    
    try:
        image = Image.open(path)
        text = pytesseract.image_to_string(image)
        
        if not text.strip():
            text = "No text detected in image."
        
        return [Document(
            page_content=text,
            metadata={"source": Path(path).name, "type": "image_ocr"}
        )]
    except Exception as e:
        return [Document(
            page_content=f"Error processing image: {str(e)}",
            metadata={"source": Path(path).name, "error": True}
        )]


def load_word_document(path: str) -> List[Document]:
    """
    Load a Word document (.docx).
    
    Args:
        path: Path to the Word document
    
    Returns:
        List containing a single Document with extracted text
    """
    if not DOCX_AVAILABLE:
        return [Document(
            page_content="Word document processing not available. Please install python-docx.",
            metadata={"source": Path(path).name, "error": True}
        )]
    
    try:
        doc = DocxDocument(path)
        full_text = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text)
                if row_text:
                    full_text.append(" | ".join(row_text))
        
        content = "\n\n".join(full_text)
        
        if not content.strip():
            content = "No text content found in document."
        
        return [Document(
            page_content=content,
            metadata={"source": Path(path).name, "type": "word"}
        )]
    except Exception as e:
        return [Document(
            page_content=f"Error processing Word document: {str(e)}",
            metadata={"source": Path(path).name, "error": True}
        )]


def load_powerpoint(path: str) -> List[Document]:
    """
    Load a PowerPoint presentation (.pptx).
    
    Args:
        path: Path to the PowerPoint file
    
    Returns:
        List of Documents, one per slide
    """
    if not PPTX_AVAILABLE:
        return [Document(
            page_content="PowerPoint processing not available. Please install python-pptx.",
            metadata={"source": Path(path).name, "error": True}
        )]
    
    try:
        prs = Presentation(path)
        documents = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                content = "\n".join(slide_text)
                documents.append(Document(
                    page_content=content,
                    metadata={
                        "source": Path(path).name,
                        "slide": slide_num,
                        "type": "powerpoint"
                    }
                ))
        
        if not documents:
            documents = [Document(
                page_content="No text content found in presentation.",
                metadata={"source": Path(path).name, "type": "powerpoint"}
            )]
        
        return documents
    except Exception as e:
        return [Document(
            page_content=f"Error processing PowerPoint: {str(e)}",
            metadata={"source": Path(path).name, "error": True}
        )]


def load_excel(path: str) -> List[Document]:
    """
    Load an Excel spreadsheet (.xlsx, .xls, .csv).
    
    Args:
        path: Path to the Excel file
    
    Returns:
        List of Documents, one per sheet
    """
    if not EXCEL_AVAILABLE:
        return [Document(
            page_content="Excel processing not available. Please install openpyxl and pandas.",
            metadata={"source": Path(path).name, "error": True}
        )]
    
    try:
        file_ext = Path(path).suffix.lower()
        
        if file_ext == '.csv':
            # Handle CSV files
            df = pd.read_csv(path)
            content = df.to_string(index=False)
            return [Document(
                page_content=content,
                metadata={"source": Path(path).name, "type": "csv"}
            )]
        else:
            # Handle Excel files
            excel_file = pd.ExcelFile(path)
            documents = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(path, sheet_name=sheet_name)
                content = f"Sheet: {sheet_name}\n\n{df.to_string(index=False)}"
                
                documents.append(Document(
                    page_content=content,
                    metadata={
                        "source": Path(path).name,
                        "sheet": sheet_name,
                        "type": "excel"
                    }
                ))
            
            return documents if documents else [Document(
                page_content="No data found in spreadsheet.",
                metadata={"source": Path(path).name, "type": "excel"}
            )]
    except Exception as e:
        return [Document(
            page_content=f"Error processing Excel file: {str(e)}",
            metadata={"source": Path(path).name, "error": True}
        )]


# ==================== Document Chunking ====================
def chunk_documents(docs: List[Document]) -> List[Document]:
    """
    Split documents into smaller chunks for better retrieval.
    
    Args:
        docs: List of Document objects to chunk
        
    Returns:
        List of chunked Document objects
    """
    # Filter out empty or error documents
    valid_docs = []
    for doc in docs:
        content = doc.page_content.strip()
        # Skip empty documents or error messages
        if content and not doc.metadata.get("error", False):
            # Ensure minimum content length
            if len(content) >= 10:
                valid_docs.append(doc)
    
    if not valid_docs:
        raise ValueError("No valid text content found in document. The file may be a scanned PDF (image-based), corrupted, or empty.")
    
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        length_function=len,
        separators=["\n\n", "\n", " ", ""]
    )
    
    chunked_docs = text_splitter.split_documents(valid_docs)
    
    # Final check to ensure we have chunks with content
    if not chunked_docs:
        raise ValueError("Document chunking produced no valid chunks. The file may contain no extractable text.")
    
    return chunked_docs


# ==================== Indexing ====================
def index_document(file_path: str, display_name: str) -> int:
    """
    Index a document into the vector store.
    
    Args:
        file_path: Path to the file to index
        display_name: Display name for the document (usually filename)
        
    Returns:
        Number of chunks added to the vector store
        
    Raises:
        ValueError: If document is empty or contains no extractable text
        Exception: For other processing errors
    """
    # Load documents
    docs = load_file_to_documents(file_path)
    
    # Check if any documents have errors
    error_docs = [doc for doc in docs if doc.metadata.get("error", False)]
    if error_docs:
        error_msg = error_docs[0].page_content
        error_type = error_docs[0].metadata.get("error_type", "unknown")
        if error_type == "scanned_pdf":
            raise ValueError("Scanned PDF detected. This file contains only images with no extractable text. Please convert to text-based PDF or use OCR.")
        else:
            raise ValueError(error_msg)
    
    # Add metadata to each document
    for doc in docs:
        doc.metadata["source"] = display_name
        # Preserve page information if it exists (from PyPDFLoader)
        if "page" not in doc.metadata and "page_number" in doc.metadata:
            doc.metadata["page"] = doc.metadata["page_number"]
    
    # Chunk documents (this will validate content and raise error if empty)
    chunked_docs = chunk_documents(docs)
    
    # Final validation before adding to vector store
    if not chunked_docs:
        raise ValueError("No valid content to index. The document may be empty or contain only images.")
    
    # Get vector store and add documents
    vector_store = get_vector_store()
    vector_store.add_documents(chunked_docs)
    
    return len(chunked_docs)


# ==================== Querying ====================
def query_notes(question: str, k: int = TOP_K) -> Dict[str, Any]:
    """
    Query the indexed notes using RAG.
    Searches across ALL documents in the library.
    
    Args:
        question: The question to answer
        k: Number of chunks to retrieve
        
    Returns:
        Dictionary containing:
        - answer: The generated answer
        - citations: List of source citations with snippets
        - sources_searched: Number of unique sources used
    """
    # Get vector store and create retriever
    vector_store = get_vector_store()
    retriever = vector_store.as_retriever(search_kwargs={"k": k})
    
    # Retrieve relevant documents
    retrieved_docs = retriever.invoke(question)
    
    # If no documents found, return early
    if not retrieved_docs:
        return {
            "answer": "I couldn't find anything related to that in your library. Please make sure you've uploaded relevant documents.",
            "citations": [],
            "sources_searched": 0
        }
    
    # Build context for the LLM
    context_parts = []
    citations = []
    sources_used = set()
    
    for i, doc in enumerate(retrieved_docs, 1):
        source = doc.metadata.get("source", "Unknown")
        sources_used.add(source)
        page = doc.metadata.get("page", None)
        slide = doc.metadata.get("slide", None)
        sheet = doc.metadata.get("sheet", None)
        doc_type = doc.metadata.get("type", "document")
        
        # Build source header with appropriate location info
        source_header = f"[{i}] Source: {source}"
        
        if page is not None:
            source_header += f" (Page {page + 1})"
        elif slide is not None:
            source_header += f" (Slide {slide})"
        elif sheet is not None:
            source_header += f" (Sheet: {sheet})"
        
        context_parts.append(f"{source_header}\n{doc.page_content}\n")
        
        # Prepare citation for response
        snippet = doc.page_content[:400] + "..." if len(doc.page_content) > 400 else doc.page_content
        citation = {
            "id": i,
            "source": source,
            "snippet": snippet,
            "type": doc_type
        }
        
        if page is not None:
            citation["page"] = page + 1
        if slide is not None:
            citation["slide"] = slide
        if sheet is not None:
            citation["sheet"] = sheet
        
        citations.append(citation)
    
    # Build the prompt
    context = "\n".join(context_parts)
    
    prompt = f"""You are a helpful study assistant with access to a comprehensive library. Answer the question based ONLY on the context provided below from multiple documents in the library.

Context from the user's library ({len(sources_used)} documents):
{context}

Instructions:
- Use ONLY the information from the context above to answer the question.
- The context comes from multiple documents in the library - synthesize information from different sources when relevant.
- If you cannot find the answer in the context, clearly state that you don't know or that the information is not in the provided documents.
- Cite your sources by using the citation numbers [1], [2], etc. when referencing specific information.
- Be clear, concise, and accurate.
- If information from multiple sources is relevant, mention all relevant sources.

Question: {question}

Answer:"""
    
    # Get LLM and generate answer
    llm = get_llm()
    response = llm.invoke(prompt)
    
    # Extract the text content from the response
    answer_text = response.content if hasattr(response, 'content') else str(response)
    
    return {
        "answer": answer_text,
        "citations": citations,
        "sources_searched": len(sources_used),
        "total_chunks": len(retrieved_docs)
    }


# ==================== Library Management ====================
def get_library_info() -> Dict[str, Any]:
    """
    Get information about all documents in the library.
    
    Returns:
        Dictionary with library statistics and file list
    """
    try:
        vector_store = get_vector_store()
        collection = vector_store._collection
        
        # Get all documents
        all_data = collection.get(include=['metadatas'])
        
        if not all_data['ids']:
            return {
                "total_documents": 0,
                "total_chunks": 0,
                "documents": []
            }
        
        # Extract unique sources
        documents_dict = {}
        for metadata in all_data['metadatas']:
            source = metadata.get('source', 'Unknown')
            if source not in documents_dict:
                documents_dict[source] = {
                    "name": source,
                    "type": metadata.get('type', 'unknown'),
                    "chunks": 0
                }
            documents_dict[source]["chunks"] += 1
        
        return {
            "total_documents": len(documents_dict),
            "total_chunks": len(all_data['ids']),
            "documents": list(documents_dict.values())
        }
    except Exception as e:
        return {
            "error": str(e),
            "total_documents": 0,
            "total_chunks": 0,
            "documents": []
        }


def delete_document_from_library(filename: str) -> Dict[str, Any]:
    """
    Delete all chunks of a document from the vector store.
    
    Args:
        filename: Name of the file to delete
        
    Returns:
        Dictionary with deletion status
    """
    try:
        vector_store = get_vector_store()
        collection = vector_store._collection
        
        # Get all document IDs with this source
        results = collection.get(
            where={"source": filename},
            include=['metadatas']
        )
        
        if not results['ids']:
            return {
                "status": "error",
                "message": f"Document '{filename}' not found in library"
            }
        
        # Delete the documents
        collection.delete(ids=results['ids'])
        
        # Also try to delete the physical file
        from pathlib import Path
        file_path = Path("./uploads") / filename
        if file_path.exists():
            file_path.unlink()
        
        return {
            "status": "success",
            "message": f"Deleted {len(results['ids'])} chunks from '{filename}'",
            "chunks_deleted": len(results['ids'])
        }
    except Exception as e:
        return {
            "status": "error",
            "message": str(e)
        }
